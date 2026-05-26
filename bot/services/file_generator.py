import io
import logging
from typing import Any, Dict, List

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PptCm
from pptx.util import Inches, Pt as PptPt

logger = logging.getLogger(__name__)

# ─── PPTX color scheme ────────────────────────────────────────────────────────

_PRESET_SCHEMES: Dict[str, Dict[str, str]] = {
    "blue": {
        "background": "#FFFFFF",
        "title_color": "#1F3864",
        "accent":      "#2E75B6",
        "text":        "#333333",
        "header_fill": "#2E75B6",
        "header_text": "#FFFFFF",
    },
    "green": {
        "background": "#FFFFFF",
        "title_color": "#1E4620",
        "accent":      "#2E7D32",
        "text":        "#333333",
        "header_fill": "#2E7D32",
        "header_text": "#FFFFFF",
    },
    "dark": {
        "background": "#1E1E2E",
        "title_color": "#FFFFFF",
        "accent":      "#7C83FD",
        "text":        "#E0E0E0",
        "header_fill": "#7C83FD",
        "header_text": "#FFFFFF",
    },
    "minimal": {
        "background": "#FAFAFA",
        "title_color": "#212121",
        "accent":      "#757575",
        "text":        "#424242",
        "header_fill": "#F5F5F5",
        "header_text": "#212121",
    },
}


def _hex_to_ppt_rgb(hex_str: str) -> PptRGB:
    h = hex_str.lstrip("#")
    try:
        return PptRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, IndexError):
        return PptRGB(0x2E, 0x75, 0xB6)


def _get_color_scheme(data: dict) -> Dict[str, PptRGB]:
    custom = data.get("custom_colors")
    if custom and isinstance(custom, dict):
        source = {
            "background": custom.get("background", "#FFFFFF"),
            "title_color": custom.get("primary", "#1F3864"),
            "accent":      custom.get("accent", "#2E75B6"),
            "text":        custom.get("text", "#333333"),
            "header_fill": custom.get("header_fill", "#2E75B6"),
            "header_text": custom.get("header_text", "#FFFFFF"),
        }
    else:
        scheme_name = data.get("color_scheme", "blue")
        source = _PRESET_SCHEMES.get(scheme_name, _PRESET_SCHEMES["blue"])
    return {k: _hex_to_ppt_rgb(v) for k, v in source.items()}


# ─── DOCX helpers ─────────────────────────────────────────────────────────────

def _set_font(run, bold: bool = False, italic: bool = False, size_pt: int = 14) -> None:
    run.font.name = "Times New Roman"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic


def _add_body_paragraph(doc: Document, text: str) -> None:
    para = doc.add_paragraph()
    para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    para.paragraph_format.first_line_indent = Cm(1.25)
    para.paragraph_format.space_after = Pt(0)
    run = para.add_run(text)
    _set_font(run)


def _add_heading(doc: Document, text: str, level: int = 1) -> None:
    para = doc.add_paragraph()
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    para.paragraph_format.space_before = Pt(12)
    para.paragraph_format.space_after = Pt(6)
    run = para.add_run(text.upper() if level == 1 else text)
    _set_font(run, bold=True, size_pt=14)


def _set_table_borders(table) -> None:
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{side}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "4")
        el.set(qn("w:space"), "0")
        el.set(qn("w:color"), "000000")
        borders.append(el)
    tbl_pr.append(borders)


def _shade_cell(cell, hex_color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tc_pr.append(shd)


def _add_table(doc: Document, table_data: dict) -> None:
    caption = table_data.get("caption", "")
    headers: List[str] = table_data.get("headers", [])
    rows: List[List[str]] = table_data.get("rows", [])

    if not headers:
        return

    if caption:
        cap_para = doc.add_paragraph()
        cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap_run = cap_para.add_run(caption)
        _set_font(cap_run, italic=True, size_pt=12)

    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    _set_table_borders(table)

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.rows[0].cells[i]
        _shade_cell(cell, "D9D9D9")
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cell.paragraphs[0].add_run(header_text)
        _set_font(run, bold=True, size_pt=12)

    # Data rows
    for r_idx, row_data in enumerate(rows, start=1):
        for c_idx, cell_value in enumerate(row_data):
            if c_idx >= len(headers):
                break
            cell = table.rows[r_idx].cells[c_idx]
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = cell.paragraphs[0].add_run(str(cell_value))
            _set_font(run, size_pt=12)

    doc.add_paragraph()  # spacing after table


def _add_footer_warning(doc: Document) -> None:
    section = doc.sections[0]
    footer = section.footer
    para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    para.clear()
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = para.add_run("⚠️ Данные замерные — замените своими")
    run.font.name = "Times New Roman"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)


# ─── DOCX public ──────────────────────────────────────────────────────────────

def create_lab_docx(data: dict, output_path: str) -> str:
    doc = Document()

    # Page setup: GOST margins
    section = doc.sections[0]
    section.left_margin = Cm(3.0)
    section.right_margin = Cm(2.0)
    section.top_margin = Cm(2.0)
    section.bottom_margin = Cm(2.0)

    # Remove default styles noise by resetting Normal style
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(14)

    # Title page block
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_para.paragraph_format.space_before = Pt(72)
    title_run = title_para.add_run(data.get("title", "Лабораторная работа"))
    _set_font(title_run, bold=True, size_pt=16)

    subject_para = doc.add_paragraph()
    subject_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subject_run = subject_para.add_run(f"Дисциплина: {data.get('subject', '')}")
    _set_font(subject_run, size_pt=14)

    if data.get("student_placeholder"):
        student_para = doc.add_paragraph()
        student_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        student_run = student_para.add_run("Выполнил(а): _______________")
        _set_font(student_run, size_pt=14)

    doc.add_page_break()

    # Sections
    for section_data in data.get("sections", []):
        _add_heading(doc, section_data.get("heading", ""))

        text = section_data.get("text", "")
        if text:
            _add_body_paragraph(doc, text)

        for formula in section_data.get("formulas", []):
            formula_para = doc.add_paragraph()
            formula_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            formula_run = formula_para.add_run(formula)
            _set_font(formula_run, italic=True)

        for table_data in section_data.get("tables", []):
            _add_table(doc, table_data)

    # Calculations
    if data.get("calculations"):
        _add_heading(doc, "Расчёты")
        _add_body_paragraph(doc, data["calculations"])

    # Conclusion
    if data.get("conclusion"):
        _add_heading(doc, "Вывод")
        _add_body_paragraph(doc, data["conclusion"])

    if data.get("has_placeholder_data"):
        _add_footer_warning(doc)

    doc.save(output_path)
    logger.info("DOCX saved: %s", output_path)
    return output_path


# ─── PPTX helpers ─────────────────────────────────────────────────────────────

_HDR_H = PptCm(1.2)
_WHITE = PptRGB(0xFF, 0xFF, 0xFF)


def _ppt_set_bg(slide, color: PptRGB) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _ppt_add_textbox(
    slide,
    text: str,
    left, top, width, height,
    font_size: int = 18,
    bold: bool = False,
    color: PptRGB | None = None,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
    font_name: str = "Calibri",
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _ppt_header_bar(slide, title: str, scheme: Dict[str, PptRGB], prs: Presentation) -> None:
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, _HDR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = scheme["header_fill"]
    bar.line.fill.background()
    _ppt_add_textbox(
        slide, title,
        PptCm(0.4), PptCm(0.15),
        prs.slide_width - PptCm(0.8), _HDR_H,
        font_size=14, bold=True,
        color=scheme["header_text"],
        align=PP_ALIGN.LEFT,
    )


def _ppt_slide_number(slide, number: int, scheme: Dict[str, PptRGB], prs: Presentation) -> None:
    w = PptCm(1.2)
    h = PptCm(0.5)
    margin = PptCm(0.3)
    _ppt_add_textbox(
        slide, str(number),
        prs.slide_width - w - margin,
        prs.slide_height - h - margin,
        w, h,
        font_size=12,
        color=scheme["accent"],
        align=PP_ALIGN.RIGHT,
    )


def _ppt_bullets(
    slide, items: List[str], scheme: Dict[str, PptRGB], prs: Presentation,
    top, height, font_size: int,
) -> None:
    content_left = PptCm(1.0)
    content_w = prs.slide_width - PptCm(2.0)
    txBox = slide.shapes.add_textbox(content_left, top, content_w, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        mr = p.add_run()
        mr.text = "▸ "
        mr.font.name = "Calibri"
        mr.font.size = PptPt(font_size)
        mr.font.color.rgb = scheme["accent"]
        tr = p.add_run()
        tr.text = str(bullet)
        tr.font.name = "Calibri"
        tr.font.size = PptPt(font_size)
        tr.font.color.rgb = scheme["text"]


def _bullet_font_size(count: int) -> int:
    if count <= 4:
        return 18
    if count <= 6:
        return 16
    if count <= 9:
        return 14
    return 12


# ─── PPTX slide builders ──────────────────────────────────────────────────────

def _ppt_build_title_slide(slide, data: dict, scheme: Dict[str, PptRGB], prs: Presentation, num: int) -> None:
    _ppt_set_bg(slide, scheme["background"])

    # Title
    _ppt_add_textbox(
        slide, data.get("title", ""),
        PptCm(1.5), PptCm(4.0),
        prs.slide_width - PptCm(3.0), PptCm(4.0),
        font_size=40, bold=True,
        color=scheme["title_color"],
        align=PP_ALIGN.CENTER,
    )

    # Subtitle
    subtitle = data.get("subtitle") or data.get("notes", "")
    if subtitle:
        _ppt_add_textbox(
            slide, subtitle,
            PptCm(2.0), PptCm(8.5),
            prs.slide_width - PptCm(4.0), PptCm(2.0),
            font_size=22,
            color=scheme["accent"],
            align=PP_ALIGN.CENTER,
        )

    # Decorative line — 40% width centered
    line_w = int(prs.slide_width * 0.4)
    line_left = (prs.slide_width - line_w) // 2
    line = slide.shapes.add_shape(1, line_left, PptCm(12.5), line_w, PptPt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = scheme["accent"]
    line.line.fill.background()

    _ppt_slide_number(slide, num, scheme, prs)


def _ppt_build_content_slide(slide, data: dict, scheme: Dict[str, PptRGB], prs: Presentation, num: int) -> None:
    _ppt_set_bg(slide, scheme["background"])
    _ppt_header_bar(slide, data.get("title", ""), scheme, prs)

    content: List[str] = data.get("content", [])
    font_size = _bullet_font_size(len(content))

    content_top = _HDR_H + PptCm(0.5)
    content_h = prs.slide_height - content_top - PptCm(0.8)
    _ppt_bullets(slide, content, scheme, prs, content_top, content_h, font_size)
    _ppt_slide_number(slide, num, scheme, prs)


def _ppt_build_two_column_slide(slide, data: dict, scheme: Dict[str, PptRGB], prs: Presentation, num: int) -> None:
    left_col: List[str] = data.get("left_column") or []
    right_col: List[str] = data.get("right_column") or []

    # Fallback: if the model returned content instead of left/right columns
    if not left_col and not right_col:
        content: List[str] = data.get("content", [])
        mid = max(1, len(content) // 2)
        left_col, right_col = content[:mid], content[mid:]

    if not left_col and not right_col:
        _ppt_build_content_slide(slide, data, scheme, prs, num)
        return

    _ppt_set_bg(slide, scheme["background"])
    _ppt_header_bar(slide, data.get("title", ""), scheme, prs)

    col_top = _HDR_H + PptCm(0.5)
    col_h = prs.slide_height - col_top - PptCm(0.8)
    font_size = _bullet_font_size(max(len(left_col), len(right_col)))
    col_w = prs.slide_width // 2 - PptCm(1.0)

    # Left column
    txl = slide.shapes.add_textbox(PptCm(0.5), col_top, col_w, col_h)
    tfl = txl.text_frame
    tfl.word_wrap = True
    for i, item in enumerate(left_col):
        p = tfl.paragraphs[0] if i == 0 else tfl.add_paragraph()
        mr = p.add_run(); mr.text = "▸ "; mr.font.name = "Calibri"; mr.font.size = PptPt(font_size); mr.font.color.rgb = scheme["accent"]
        tr = p.add_run(); tr.text = str(item); tr.font.name = "Calibri"; tr.font.size = PptPt(font_size); tr.font.color.rgb = scheme["text"]

    # Vertical divider
    div_x = prs.slide_width // 2
    div = slide.shapes.add_shape(1, div_x, col_top, PptPt(1), col_h)
    div.fill.solid(); div.fill.fore_color.rgb = scheme["accent"]; div.line.fill.background()

    # Right column
    txr = slide.shapes.add_textbox(div_x + PptCm(0.3), col_top, col_w, col_h)
    tfr = txr.text_frame
    tfr.word_wrap = True
    for i, item in enumerate(right_col):
        p = tfr.paragraphs[0] if i == 0 else tfr.add_paragraph()
        mr = p.add_run(); mr.text = "▸ "; mr.font.name = "Calibri"; mr.font.size = PptPt(font_size); mr.font.color.rgb = scheme["accent"]
        tr = p.add_run(); tr.text = str(item); tr.font.name = "Calibri"; tr.font.size = PptPt(font_size); tr.font.color.rgb = scheme["text"]

    _ppt_slide_number(slide, num, scheme, prs)


def _ppt_build_section_header_slide(slide, data: dict, scheme: Dict[str, PptRGB], prs: Presentation, num: int) -> None:
    _ppt_set_bg(slide, scheme["background"])

    # Colored band in the center
    band_h = PptCm(5.0)
    band_top = (prs.slide_height - band_h) // 2
    band = slide.shapes.add_shape(1, 0, band_top, prs.slide_width, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = scheme["header_fill"]
    band.line.fill.background()

    _ppt_add_textbox(
        slide, data.get("title", ""),
        PptCm(1.0), band_top + PptCm(0.5),
        prs.slide_width - PptCm(2.0), band_h - PptCm(1.0),
        font_size=36, bold=True,
        color=scheme["header_text"],
        align=PP_ALIGN.CENTER,
    )
    _ppt_slide_number(slide, num, scheme, prs)


def _ppt_build_conclusion_slide(slide, data: dict, scheme: Dict[str, PptRGB], prs: Presentation, num: int) -> None:
    _ppt_set_bg(slide, scheme["background"])
    _ppt_header_bar(slide, data.get("title", ""), scheme, prs)

    content: List[str] = data.get("content", [])
    font_size = _bullet_font_size(len(content))
    footer_h = PptCm(1.5)
    content_top = _HDR_H + PptCm(0.5)
    content_h = prs.slide_height - content_top - footer_h - PptCm(0.3)
    _ppt_bullets(slide, content, scheme, prs, content_top, content_h, font_size)

    # Footer bar
    footer_top = prs.slide_height - footer_h
    bar = slide.shapes.add_shape(1, 0, footer_top, prs.slide_width, footer_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = scheme["accent"]
    bar.line.fill.background()
    _ppt_add_textbox(
        slide, "Спасибо за внимание!",
        PptCm(0.5), footer_top + PptCm(0.2),
        prs.slide_width - PptCm(1.0), footer_h,
        font_size=20, bold=True,
        color=_WHITE,
        align=PP_ALIGN.CENTER,
    )
    _ppt_slide_number(slide, num, scheme, prs)


_SLIDE_BUILDERS = {
    "title":          _ppt_build_title_slide,
    "content":        _ppt_build_content_slide,
    "two_column":     _ppt_build_two_column_slide,
    "section_header": _ppt_build_section_header_slide,
    "conclusion":     _ppt_build_conclusion_slide,
    "image_text":     _ppt_build_content_slide,  # graceful fallback
}


# ─── PPTX public ──────────────────────────────────────────────────────────────

def create_presentation_pptx(data: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # 33.87 cm — 16:9 widescreen
    prs.slide_height = Inches(7.5)    # 19.05 cm

    blank_layout = prs.slide_layouts[6]  # Blank
    scheme = _get_color_scheme(data)

    slides_data = data.get("slides", [])
    if not slides_data:
        slides_data = [{"slide_number": 1, "layout": "title", "title": data.get("title", ""), "content": []}]

    for i, slide_data in enumerate(slides_data, start=1):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_data.get("layout", "content")
        builder = _SLIDE_BUILDERS.get(layout, _ppt_build_content_slide)
        try:
            builder(slide, slide_data, scheme, prs, i)
        except Exception as e:
            logger.error("Slide build error slide=%s layout=%s: %s", i, layout, e)
            _ppt_build_content_slide(slide, slide_data, scheme, prs, i)

    prs.save(output_path)
    logger.info("PPTX saved: %s", output_path)
    return output_path


# ─── Text formatter ──────────────────────────────────────────────────────────

def format_text_answer(data: dict) -> List[str]:
    """Splits answer into Telegram-safe chunks (≤4096 chars)."""
    answer: str = data.get("answer", "")
    max_len = 4096

    if len(answer) <= max_len:
        return [answer]

    parts: List[str] = []
    while answer:
        if len(answer) <= max_len:
            parts.append(answer)
            break

        # Split at last newline within limit to avoid cutting mid-sentence
        cut = answer.rfind("\n", 0, max_len)
        if cut == -1:
            cut = max_len

        parts.append(answer[:cut])
        answer = answer[cut:].lstrip("\n")

    return parts
